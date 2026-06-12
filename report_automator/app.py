"""
app.py - Entry point aplikasi Streamlit untuk Automasi Laporan Telkomsel
=========================================================================
Template: Coverage Activity Report
Struktur:
  Slide 1       : Cover
  Slide 2..N    : Satu slide per baris data (Site)
  Slide terakhir: Penutup

Mode:
  - Proposal Only  : file Excel dengan sheet Proposal saja (atau 1 sheet)
  - Proposal+Result: file Excel dengan 2 sheet (Proposal + Result)
"""

import streamlit as st
import pandas as pd
import io
from datetime import datetime

from data_processor import (
    get_excel_info, load_excel_sheet, load_data,
    detect_sheet_mode, load_result_sheet,
    COL_RESULT_BG, COL_RESULT_PR,
)
from ppt_generator import (
    generate_pptx,
    COL_SITE_ID, COL_SITE_NAME, COL_PURPOSE, COL_CITY,
    COL_FINDING, COL_PLAN_ACTION, COL_SUPPORT, COL_GOALS, COL_INCREMENT,
)
from ppt_generator_result import generate_result_pptx


# ─────────────────────────────────────────────────────────────────────────────
# HELPER: Merge dua file PPTX menjadi satu
# ─────────────────────────────────────────────────────────────────────────────
def merge_pptx(buf_a: bytes, buf_b: bytes) -> bytes:
    """
    Gabungkan dua file PPTX (Proposal + Result) menjadi satu file.

    Struktur hasil gabungan:
    ┌─────────────────────────────────────────────────┐
    │  Slide 1        : Cover (dari Proposal)          │
    │  Slide 2..N     : Isi Proposal (tanpa cover &    │
    │                   tanpa penutup)                 │
    │  Slide N+1..M   : Isi Result (tanpa cover &      │
    │                   tanpa penutup)                 │
    │  Slide terakhir : Penutup (dari Result)          │
    └─────────────────────────────────────────────────┘

    Best practice:
    - Tidak mengubah slide asli (non-destructive)
    - Menggunakan copy.deepcopy agar layout/theme tidak tertimpa
    - Cover hanya 1x di awal, Penutup hanya 1x di akhir
    """
    import copy
    from pptx import Presentation
    from lxml import etree

    # Buka kedua presentasi dari bytes
    prs_a = Presentation(io.BytesIO(buf_a))  # Proposal (akan jadi base)
    prs_b = Presentation(io.BytesIO(buf_b))  # Result (akan di-append)

    def _clone_slide(prs_target, slide_src, layout):
        """Clone satu slide dari sumber ke presentasi target."""
        new_slide = prs_target.slides.add_slide(layout)
        sp_tree = new_slide.shapes._spTree
        for child in list(sp_tree):
            sp_tree.remove(child)
        src_sp_tree = slide_src.shapes._spTree
        for child in src_sp_tree:
            sp_tree.append(copy.deepcopy(child))
        for rel in slide_src.part.rels.values():
            if "image" in rel.reltype:
                try:
                    img_part = rel.target_part
                    new_slide.part.relate_to(img_part, rel.reltype)
                except Exception:
                    pass

    # Referensi layout blank
    blank_layout = prs_a.slide_layouts[6]

    slides_a = list(prs_a.slides)
    slides_b = list(prs_b.slides)

    # ── Langkah 1: Hapus slide PENUTUP dari Proposal (slide terakhir prs_a) ───
    # Kita rebuild: ambil semua slide A kecuali terakhir, lalu tambahkan
    # slide-slide Result (kecuali cover & penutupnya), lalu tambahkan penutup.
    #
    # python-pptx tidak mendukung penghapusan slide langsung, jadi kita
    # buat presentasi baru dan clone ulang slide yang dibutuhkan.

    prs_out = Presentation()
    prs_out.slide_width  = prs_a.slide_width
    prs_out.slide_height = prs_a.slide_height
    blank_out = prs_out.slide_layouts[6]

    def _clone_to(prs_target, slide_src, layout):
        """Clone satu slide ke presentasi target."""
        new_slide = prs_target.slides.add_slide(layout)
        sp_tree = new_slide.shapes._spTree
        for child in list(sp_tree):
            sp_tree.remove(child)
        for child in slide_src.shapes._spTree:
            sp_tree.append(copy.deepcopy(child))
        for rel in slide_src.part.rels.values():
            if "image" in rel.reltype:
                try:
                    prs_target.slides[-1].part.relate_to(
                        rel.target_part, rel.reltype
                    )
                except Exception:
                    pass

    # Clone semua slide Proposal KECUALI slide terakhir (Penutup Proposal)
    proposal_content = slides_a[:-1] if len(slides_a) > 1 else slides_a
    for sl in proposal_content:
        _clone_to(prs_out, sl, blank_out)

    # Clone slide Result: skip slide pertama (Cover) dan slide terakhir (Penutup Result)
    # Simpan slide penutup Result untuk ditambahkan di akhir
    if len(slides_b) > 2:
        result_content = slides_b[1:-1]   # isi materi Result
        result_closing = slides_b[-1]     # Penutup Result
    elif len(slides_b) == 2:
        result_content = []               # hanya Cover + Penutup, tidak ada isi
        result_closing = slides_b[-1]
    else:
        result_content = []
        result_closing = slides_b[-1] if slides_b else None

    for sl in result_content:
        _clone_to(prs_out, sl, blank_out)

    # Tambahkan Penutup (dari Result) sebagai slide terakhir
    if result_closing is not None:
        _clone_to(prs_out, result_closing, blank_out)

    # Simpan ke buffer bytes
    out_buf = io.BytesIO()
    prs_out.save(out_buf)
    out_buf.seek(0)
    return out_buf.getvalue()


# ─────────────────────────────────────────────────────────────────────────────
# Konfigurasi halaman
# ─────────────────────────────────────────────────────────────────────────────
st.set_page_config(
    page_title=" Telkomsel Report Automator",
    page_icon="📡",
    layout="wide",
    initial_sidebar_state="expanded",
)

st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap');
    html, body, [class*="css"] { font-family: 'Inter', sans-serif; }

    .main-header {
        background: linear-gradient(135deg, #8B0000, #CC0000, #FF3333);
        color: white;
        padding: 1.8rem 2.5rem;
        border-radius: 16px;
        margin-bottom: 1.5rem;
        box-shadow: 0 8px 32px rgba(204,0,0,0.35);
    }
    .main-header h1 { font-size: 2rem; font-weight: 700; margin: 0; }
    .main-header p  { margin: 0.4rem 0 0 0; opacity: 0.85; font-size: 0.95rem; }

    .step-badge {
        background: linear-gradient(135deg, #CC0000, #FF3333);
        color: white;
        padding: 0.35rem 1rem;
        border-radius: 20px;
        font-weight: 700;
        font-size: 0.8rem;
        display: inline-block;
        margin-bottom: 0.5rem;
    }

    .mode-badge-proposal {
        background: linear-gradient(135deg, #1B2A4A, #2d4270);
        color: white;
        padding: 0.5rem 1.2rem;
        border-radius: 10px;
        font-weight: 600;
        font-size: 0.85rem;
        border-left: 4px solid #CC0000;
        margin-bottom: 1rem;
        display: inline-block;
    }
    .mode-badge-result {
        background: linear-gradient(135deg, #004d00, #006600);
        color: white;
        padding: 0.5rem 1.2rem;
        border-radius: 10px;
        font-weight: 600;
        font-size: 0.85rem;
        border-left: 4px solid #00CC44;
        margin-bottom: 1rem;
        display: inline-block;
    }

    .site-card {
        background: #1a1a2e;
        border: 1px solid #CC0000;
        border-radius: 10px;
        padding: 1rem 1.2rem;
        margin-bottom: 0.8rem;
    }

    section[data-testid="stSidebar"] { background: #0d1117; }
    section[data-testid="stSidebar"] * { color: #e6edf3 !important; }

    .stDownloadButton > button {
        background: linear-gradient(135deg, #CC0000, #990000);
        color: white; border: none; border-radius: 10px;
        padding: 0.7rem 2rem; font-weight: 700; font-size: 1rem;
        width: 100%; transition: all 0.3s ease;
        box-shadow: 0 4px 15px rgba(204,0,0,0.4);
    }
    .stDownloadButton > button:hover {
        transform: translateY(-2px);
        box-shadow: 0 6px 20px rgba(204,0,0,0.6);
    }
    .section-divider {
        border: none; height: 1px;
        background: linear-gradient(to right, transparent, #CC0000, transparent);
        margin: 1.5rem 0;
    }
    .result-section-divider {
        border: none; height: 1px;
        background: linear-gradient(to right, transparent, #00AA44, transparent);
        margin: 1.5rem 0;
    }

    /* Tab styling */
    .stTabs [data-baseweb="tab-list"] {
        gap: 8px;
    }
    .stTabs [data-baseweb="tab"] {
        border-radius: 8px 8px 0 0;
        padding: 0.5rem 1.2rem;
        font-weight: 600;
    }
</style>
""", unsafe_allow_html=True)


# ─────────────────────────────────────────────────────────────────────────────
# Header
# ─────────────────────────────────────────────────────────────────────────────
st.markdown("""
<div class="main-header">
    <h1>📡 Telkomsel Coverage Report Automator</h1>
    <p>Upload data Excel → Preview per Site → Generate PowerPoint otomatis (Proposal & Result)</p>
</div>
""", unsafe_allow_html=True)


# ─────────────────────────────────────────────────────────────────────────────
# Sidebar: Konfigurasi Laporan
# ─────────────────────────────────────────────────────────────────────────────
with st.sidebar:
    st.markdown("##  Konfigurasi Laporan")
    st.markdown("---")

    report_title = st.text_input(
        " Judul Laporan",
        value="Tracking Activity NOP 2026",
        help="Judul ini akan muncul di Slide Cover.",
    )
    report_subtitle = st.text_input(
        " Sub-judul / Departemen",
        value="Divisi Network Operation",
        help="Sub-judul di bawah judul utama.",
    )
    author_name = st.text_input(
        " Nama Pembuat",
        value="Tim Network Operation",
    )
    report_date = st.date_input(
        " Tanggal Laporan",
        value=datetime.today(),
    )

    st.markdown("---")
    st.markdown(
        "<small> Data hanya diproses lokal dan tidak dikirim ke server manapun.</small>",
        unsafe_allow_html=True,
    )


# ─────────────────────────────────────────────────────────────────────────────
# Step indicator
# ─────────────────────────────────────────────────────────────────────────────
c1, c2, c3 = st.columns(3)
with c1: st.info("**Step 1** · Upload Excel Data")
with c2: st.info("**Step 2** · Upload Gambar per Site")
with c3: st.info("**Step 3** · Generate & Download .pptx")

st.markdown('<hr class="section-divider">', unsafe_allow_html=True)


# ─────────────────────────────────────────────────────────────────────────────
# STEP 1: Upload File
# ─────────────────────────────────────────────────────────────────────────────
st.markdown("###  Step 1 · Upload File Data")
st.markdown(
    "Upload file Excel dengan sheet **Proposal** (wajib) dan opsional sheet **Result**. "
    "Sistem akan otomatis mendeteksi apakah file memiliki kedua sheet."
)

uploaded_file = st.file_uploader(
    "Seret & lepas file Excel di sini, atau klik untuk memilih",
    type=["csv", "xlsx", "xls"],
    help="Format: CSV (.csv) atau Excel (.xlsx/.xls). Setiap baris = 1 slide PPT.",
)

df_proposal = None
df_result   = None
sheet_mode  = None   # dict dari detect_sheet_mode
file_bytes  = None

if uploaded_file is not None:
    filename_lower = uploaded_file.name.lower()
    is_excel = filename_lower.endswith((".xlsx", ".xls"))

    file_bytes = uploaded_file.getvalue()

    if is_excel:
        # ── Deteksi mode sheet ────────────────────────────────────────────
        try:
            sheet_mode = detect_sheet_mode(file_bytes)
        except Exception as e:
            st.error(f"! Gagal membaca struktur file Excel: {e}")
            st.stop()

        # ── Badge mode ────────────────────────────────────────────────────
        if sheet_mode["has_result"]:
            st.markdown(
                f'<div class="mode-badge-result"> Terdeteksi 2 Sheet: '
                f'<strong>{sheet_mode["proposal_sheet"]}</strong> (Proposal) + '
                f'<strong>{sheet_mode["result_sheet"]}</strong> (Result)</div>',
                unsafe_allow_html=True
            )
        else:
            st.markdown(
                f'<div class="mode-badge-proposal"> Terdeteksi 1 Sheet: '
                f'<strong>{sheet_mode["proposal_sheet"]}</strong> (Proposal only)</div>',
                unsafe_allow_html=True
            )

        # ── Baca sheet Proposal ───────────────────────────────────────────
        try:
            df_proposal = load_excel_sheet(file_bytes, sheet_mode["proposal_sheet"])
        except Exception as e:
            st.error(f"! Gagal membaca sheet Proposal: {e}")
            st.stop()

        # ── Baca sheet Result (jika ada) ───────────────────────────────────
        if sheet_mode["has_result"]:
            try:
                df_result = load_result_sheet(file_bytes, sheet_mode["result_sheet"])
            except Exception as e:
                st.warning(f"! Sheet Result ditemukan tapi gagal dibaca: {e}")
                df_result = None

    else:
        # CSV: langsung load sebagai Proposal
        try:
            df_proposal = load_data(uploaded_file)
        except Exception as e:
            st.error(f"! Gagal membaca file: {e}")
            st.stop()

    # ── Info sukses ───────────────────────────────────────────────────────
    if df_proposal is not None:
        st.success(
            f" **{uploaded_file.name}** berhasil dimuat — "
            f"Sheet Proposal: **{df_proposal.shape[0]:,} baris** × **{df_proposal.shape[1]} kolom**"
            + (f"  |  Sheet Result: **{df_result.shape[0]:,} baris**"
               if df_result is not None else "")
        )

        col_info1, col_info2 = st.columns(2)
        with col_info1:
            st.markdown(
                f" Proposal: **{df_proposal.shape[0]} slide data** + Cover + Penutup "
                f"= **{df_proposal.shape[0] + 2} slide**"
            )
        if df_result is not None:
            with col_info2:
                st.markdown(
                    f" Result: **{df_result.shape[0]} slide data** + Cover + Penutup "
                    f"= **{df_result.shape[0] + 2} slide**"
                )

    # Preview tabel ringkas (Proposal)
    with st.expander(" Preview data Proposal (10 baris pertama)", expanded=False):
        df_preview = df_proposal.head(10).astype(str)
        st.dataframe(df_preview, use_container_width=True)

    if df_result is not None:
        with st.expander(" Preview data Result (10 baris pertama)", expanded=False):
            df_result_preview = df_result.head(10).astype(str)
            st.dataframe(df_result_preview, use_container_width=True)

    st.markdown('<hr class="section-divider">', unsafe_allow_html=True)

    # ─────────────────────────────────────────────────────────────────────────
    # STEP 2: Upload Gambar per Site — Tabs: Proposal | Result
    # ─────────────────────────────────────────────────────────────────────────
    st.markdown("###  Step 2 · Upload Gambar per Site")

    # Siapkan tab (tampilkan tab Result hanya jika data Result tersedia)
    if df_result is not None:
        tab_proposal, tab_result = st.tabs([" Gambar Proposal", " Gambar Result"])
    else:
        tab_proposal = st.tabs([" Gambar Proposal"])[0]
        tab_result   = None

    # ═══════════════════════════════════════════════════════════════════════
    # TAB PROPOSAL: Upload gambar (flow yang sudah ada)
    # ═══════════════════════════════════════════════════════════════════════
    with tab_proposal:
        st.markdown(
            "Untuk setiap site, upload gambar pendukung. "
            "Setiap site membutuhkan **minimum 3 gambar** (Payload Chart, Maps/COVMO, Support Image), "
            "dan **maksimal 5 gambar** (Payload Chart, Maps/COVMO, + 3 Support Images)."
        )

        col_site_id = COL_SITE_ID
        if col_site_id not in df_proposal.columns:
            candidates = [c for c in df_proposal.columns if "site" in c.lower() and "id" in c.lower()]
            col_site_id = candidates[0] if candidates else df_proposal.columns[0]
            st.warning(f" Kolom '{COL_SITE_ID}' tidak ditemukan. Menggunakan kolom '{col_site_id}'.")

        raw_ids = df_proposal[col_site_id].tolist()
        site_ids = []
        seen = {}
        for i, sid_raw in enumerate(raw_ids):
            sid = str(sid_raw).strip() if sid_raw is not None else ""
            if not sid or sid.lower() in ("nan", "none", ""):
                sid = f"site_{i+1}"
            if sid in seen:
                seen[sid] += 1
                sid = f"{sid}_{seen[sid]}"
            else:
                seen[sid] = 0
            site_ids.append(sid)

        site_images: dict = {}

        upload_mode = st.radio(
            "Mode upload gambar Proposal:",
            ["Upload gambar per site (direkomendasikan)", "Skip — generate tanpa gambar (placeholder)"],
            index=0,
            horizontal=True,
            key="upload_mode_proposal",
        )

        if upload_mode.startswith("Upload"):
            st.markdown("---")
            selected_sites = st.multiselect(
                "Pilih site yang ingin diisi gambarnya (kosongkan untuk semua):",
                options=site_ids,
                default=[],
                help="Jika dikosongkan, form upload ditampilkan untuk 10 site pertama.",
                key="sel_sites_proposal",
            )
            sites_to_show = selected_sites if selected_sites else site_ids[:10]

            if not selected_sites and len(site_ids) > 10:
                st.info(
                    f"! Menampilkan 10 site pertama dari {len(site_ids)} site. "
                    "Gunakan filter di atas untuk site tertentu."
                )

            for i, sid in enumerate(sites_to_show):
                orig_idx = site_ids.index(sid)
                site_name_col = COL_SITE_NAME
                if site_name_col in df_proposal.columns:
                    sname = str(df_proposal.iloc[orig_idx][site_name_col]).strip()
                    site_label = f"{sid} — {sname}" if sname else sid
                else:
                    site_label = str(sid)

                with st.expander(f" {site_label}", expanded=False):
                    col_a, col_b = st.columns([1, 1])

                    with col_a:
                        st.markdown("**① Payload Site Surrounding** *(wajib)*")
                        img_payload = st.file_uploader(
                            "Upload chart payload",
                            type=["png", "jpg", "jpeg"],
                            key=f"payload_{sid}",
                        )

                        st.markdown("**② Maps Preview & COVMO** *(wajib)*")
                        img_maps = st.file_uploader(
                            "Upload maps/COVMO",
                            type=["png", "jpg", "jpeg"],
                            key=f"maps_{sid}",
                        )

                    with col_b:
                        st.markdown("**③–⑤ Support Images** *(min 1, maks 3)*")
                        img_supports = st.file_uploader(
                            "Upload support images (1–3 gambar)",
                            type=["png", "jpg", "jpeg"],
                            accept_multiple_files=True,
                            key=f"support_{sid}",
                        )

                    imgs_this_site = []
                    if img_payload:
                        imgs_this_site.append({"bytes": img_payload.getvalue(), "label": "Payload Site"})
                    if img_maps:
                        imgs_this_site.append({"bytes": img_maps.getvalue(), "label": "Maps & COVMO"})
                    for si_img in (img_supports or [])[:3]:
                        imgs_this_site.append({"bytes": si_img.getvalue(), "label": "Support"})

                    if imgs_this_site:
                        site_images[sid] = imgs_this_site
                        st.success(f" {len(imgs_this_site)} gambar siap untuk site ini.")
                    else:
                        st.caption("Belum ada gambar — akan menggunakan placeholder.")

    # ═══════════════════════════════════════════════════════════════════════
    # TAB RESULT: Upload gambar (5 kategori, 10 slot)
    # ═══════════════════════════════════════════════════════════════════════
    site_result_images: dict = {}   # key = str(index), value = dict slot

    if tab_result is not None and df_result is not None:
        with tab_result:
            st.markdown(
                "Upload gambar untuk setiap baris di sheet **Result**. "
                "Setiap baris membutuhkan **10 capture** yang terbagi dalam 5 kategori:"
            )

            # Info card 3 section
            ic1, ic2, ic3 = st.columns(3)
            with ic1:
                st.info("**Section 1 · Background**\n\n"
                        "> Site Mapping (1)\n\n"
                        "> RSRP & RSRQ Before (1)")
            with ic2:
                st.info("**Section 2 · Experience & Docs**\n\n"
                        "> Before (3 capture)\n\n"
                        "> After (3 capture)")
            with ic3:
                st.info("**Section 3 · Productivity Result**\n\n"
                        "> RSRP & RSRQ After (2 capture)")

            st.markdown("---")

            # Tentukan jumlah baris Result
            n_result = df_result.shape[0]

            # Buat daftar label per baris Result
            # Ambil label dari Proposal (PURPOSE HEADER) jika tersedia (index-based)
            result_labels = []
            proposal_rows_list = df_proposal.to_dict(orient="records") if df_proposal is not None else []
            for r_idx in range(n_result):
                if r_idx < len(proposal_rows_list):
                    ph = str(proposal_rows_list[r_idx].get(COL_PURPOSE, "")).strip()
                    sow = str(proposal_rows_list[r_idx].get("SOW", "")).strip()
                    lbl = ph or sow or f"Site {r_idx + 1}"
                else:
                    lbl = f"Site {r_idx + 1}"
                result_labels.append(lbl)

            # Filter: pilih site mana yang ingin diisi
            selected_result_indices = st.multiselect(
                "Pilih site Result yang ingin diisi gambarnya (kosongkan untuk semua):",
                options=list(range(n_result)),
                format_func=lambda i: f"[{i+1}] {result_labels[i]}",
                default=[],
                help="Jika dikosongkan, 10 pertama ditampilkan.",
                key="sel_sites_result",
            )
            indices_to_show = (
                selected_result_indices if selected_result_indices
                else list(range(min(10, n_result)))
            )

            if not selected_result_indices and n_result > 10:
                st.info(
                    f"! Menampilkan 10 baris pertama dari {n_result} baris Result. "
                    "Gunakan filter di atas untuk baris tertentu."
                )

            for r_idx in indices_to_show:
                lbl = result_labels[r_idx]

                with st.expander(f" [{r_idx+1}] {lbl}", expanded=False):
                    # 5 tab inner per site
                    t_sm, t_rb, t_bef, t_aft, t_ra = st.tabs([
                        " Site Mapping",
                        " RSRP Before",
                        " Before (3)",
                        " After (3)",
                        " RSRP After (2)",
                    ])

                    slot_data: dict = {}

                    with t_sm:
                        st.markdown("**Site Mapping** · 1 capture")
                        f_sm = st.file_uploader(
                            "Upload Site Mapping",
                            type=["png", "jpg", "jpeg"],
                            key=f"res_sm_{r_idx}",
                        )
                        if f_sm:
                            slot_data["site_mapping"] = [f_sm.getvalue()]
                            st.image(f_sm, width=200)

                    with t_rb:
                        st.markdown("**RSRP & RSRQ (Before)** · 1 capture")
                        f_rb = st.file_uploader(
                            "Upload RSRP Before",
                            type=["png", "jpg", "jpeg"],
                            key=f"res_rb_{r_idx}",
                        )
                        if f_rb:
                            slot_data["rsrp_before"] = [f_rb.getvalue()]
                            st.image(f_rb, width=200)

                    with t_bef:
                        st.markdown("**Before Documentation** · 3 capture")
                        f_bef = st.file_uploader(
                            "Upload Before (1–3 gambar)",
                            type=["png", "jpg", "jpeg"],
                            accept_multiple_files=True,
                            key=f"res_bef_{r_idx}",
                        )
                        if f_bef:
                            before_bytes = [f.getvalue() for f in f_bef[:3]]
                            slot_data["before"] = before_bytes
                            cols_prev = st.columns(len(before_bytes))
                            for ci, fb in enumerate(f_bef[:3]):
                                cols_prev[ci].image(fb, width=130)

                    with t_aft:
                        st.markdown("**After Documentation** · 3 capture")
                        f_aft = st.file_uploader(
                            "Upload After (1–3 gambar)",
                            type=["png", "jpg", "jpeg"],
                            accept_multiple_files=True,
                            key=f"res_aft_{r_idx}",
                        )
                        if f_aft:
                            after_bytes = [f.getvalue() for f in f_aft[:3]]
                            slot_data["after"] = after_bytes
                            cols_prev2 = st.columns(len(after_bytes))
                            for ci, fa in enumerate(f_aft[:3]):
                                cols_prev2[ci].image(fa, width=130)

                    with t_ra:
                        st.markdown("**RSRP & RSRQ (After)** · 2 capture")
                        f_ra = st.file_uploader(
                            "Upload RSRP After (1–2 gambar)",
                            type=["png", "jpg", "jpeg"],
                            accept_multiple_files=True,
                            key=f"res_ra_{r_idx}",
                        )
                        if f_ra:
                            ra_bytes = [f.getvalue() for f in f_ra[:2]]
                            slot_data["rsrp_after"] = ra_bytes
                            cols_prev3 = st.columns(len(ra_bytes))
                            for ci, fra in enumerate(f_ra[:2]):
                                cols_prev3[ci].image(fra, width=150)

                    # Simpan ke dict utama
                    if slot_data:
                        site_result_images[str(r_idx)] = slot_data
                        filled_slots = list(slot_data.keys())
                        total_imgs = sum(len(v) for v in slot_data.values())
                        st.success(
                            f" {total_imgs} gambar di {len(filled_slots)} kategori siap: "
                            f"{', '.join(filled_slots)}"
                        )
                    else:
                        st.caption("Belum ada gambar — akan menggunakan placeholder.")

    st.markdown('<hr class="section-divider">', unsafe_allow_html=True)

    # ─────────────────────────────────────────────────────────────────────────
    # STEP 3: Generate & Download
    # ─────────────────────────────────────────────────────────────────────────
    st.markdown("###  Step 3 · Generate & Download Laporan PowerPoint")

    # ── Dua kolom: Proposal (kiri) dan Result (kanan, opsional) ────────────
    if df_result is not None:
        gen_col1, gen_col2 = st.columns(2)
    else:
        gen_col1 = st.container()
        gen_col2 = None

    # ── Generate Proposal ─────────────────────────────────────────────────
    with gen_col1:
        st.markdown("#### · Laporan Proposal")
        if st.button(" Generate File Proposal .pptx",
                     type="primary", use_container_width=True,
                     key="btn_gen_proposal"):
            with st.spinner("Membuat PPT Proposal..."):
                try:
                    pptx_buf = generate_pptx(
                        title=report_title,
                        subtitle=report_subtitle,
                        author=author_name,
                        report_date=report_date.strftime("%d %B %Y"),
                        df_raw=df_proposal,
                        site_images=site_images if site_images else {},
                    )
                    st.session_state["pptx_proposal_buf"] = pptx_buf
                    st.session_state["pptx_proposal_fname"] = (
                        f"Proposal_{report_title.replace(' ', '_')}_"
                        f"{report_date.strftime('%Y%m%d')}.pptx"
                    )
                    st.success(
                        f" PPT Proposal berhasil! "
                        f"Total **{df_proposal.shape[0] + 2} slide**."
                    )
                except Exception as e:
                    st.error(f"! Gagal membuat PPT Proposal: {e}")
                    import traceback
                    st.code(traceback.format_exc())

        if "pptx_proposal_buf" in st.session_state:
            st.download_button(
                label="⬇ Download Proposal .pptx",
                data=st.session_state["pptx_proposal_buf"],
                file_name=st.session_state["pptx_proposal_fname"],
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
                key="dl_proposal",
            )

    # ── Generate Result ───────────────────────────────────────────────────
    if gen_col2 is not None and df_result is not None:
        with gen_col2:
            st.markdown("#### · Laporan Result")
            if st.button(" Generate File Result .pptx",
                         type="primary", use_container_width=True,
                         key="btn_gen_result"):
                with st.spinner("Membuat PPT Result..."):
                    try:
                        pptx_result_buf = generate_result_pptx(
                            title=report_title,
                            subtitle=report_subtitle,
                            author=author_name,
                            report_date=report_date.strftime("%d %B %Y"),
                            df_proposal=df_proposal,
                            df_result=df_result,
                            site_result_images=site_result_images if site_result_images else {},
                        )
                        st.session_state["pptx_result_buf"] = pptx_result_buf
                        st.session_state["pptx_result_fname"] = (
                            f"Result_{report_title.replace(' ', '_')}_"
                            f"{report_date.strftime('%Y%m%d')}.pptx"
                        )
                        st.success(
                            f" PPT Result berhasil! "
                            f"Total **{df_result.shape[0] + 2} slide**."
                        )
                    except Exception as e:
                        st.error(f"! Gagal membuat PPT Result: {e}")
                        import traceback
                        st.code(traceback.format_exc())

            if "pptx_result_buf" in st.session_state:
                st.download_button(
                    label="⬇ Download Result .pptx",
                    data=st.session_state["pptx_result_buf"],
                    file_name=st.session_state["pptx_result_fname"],
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True,
                    key="dl_result",
                )

    # ── Download Gabungan (Proposal + Result) — selalu tampil jika ada Result
    if df_result is not None:
        st.divider()
        st.markdown("####  · Download Gabungan (Proposal + Result)")
        st.caption(
            "Generate **keduanya sekaligus** dalam satu klik, atau download "
            "masing-masing secara terpisah dari kolom di atas."
        )

        # Tombol "Generate Semua" — otomatis generate Proposal + Result + merge
        if st.button(
            " Generate Semua & Gabungkan (Proposal + Result)",
            type="primary",
            use_container_width=True,
            key="btn_gen_all",
        ):
            combined_fname = (
                f"Combined_{report_title.replace(' ', '_')}_"
                f"{report_date.strftime('%Y%m%d')}.pptx"
            )
            try:
                # 1) Generate Proposal
                with st.spinner("⏳ Membuat PPT Proposal..."):
                    pptx_buf = generate_pptx(
                        title=report_title,
                        subtitle=report_subtitle,
                        author=author_name,
                        report_date=report_date.strftime("%d %B %Y"),
                        df_raw=df_proposal,
                        site_images=site_images if site_images else {},
                    )
                    st.session_state["pptx_proposal_buf"] = pptx_buf
                    st.session_state["pptx_proposal_fname"] = (
                        f"Proposal_{report_title.replace(' ', '_')}_"
                        f"{report_date.strftime('%Y%m%d')}.pptx"
                    )

                # 2) Generate Result
                with st.spinner("⏳ Membuat PPT Result..."):
                    pptx_result_buf = generate_result_pptx(
                        title=report_title,
                        subtitle=report_subtitle,
                        author=author_name,
                        report_date=report_date.strftime("%d %B %Y"),
                        df_proposal=df_proposal,
                        df_result=df_result,
                        site_result_images=site_result_images if site_result_images else {},
                    )
                    st.session_state["pptx_result_buf"] = pptx_result_buf
                    st.session_state["pptx_result_fname"] = (
                        f"Result_{report_title.replace(' ', '_')}_"
                        f"{report_date.strftime('%Y%m%d')}.pptx"
                    )

                # 3) Merge keduanya
                with st.spinner("⏳ Menggabungkan..."):
                    combined_buf = merge_pptx(pptx_buf, pptx_result_buf)
                    st.session_state["pptx_combined_buf"] = combined_buf
                    st.session_state["pptx_combined_fname"] = combined_fname

                n_proposal_data = df_proposal.shape[0]   # hanya slide data
                n_result_data   = df_result.shape[0]      # hanya slide data
                # Gabungan: 1 Cover + data Proposal + data Result + 1 Penutup
                n_combined = 1 + n_proposal_data + n_result_data + 1
                st.success(
                    f"✅ Semua berhasil dibuat! "
                    f"Total **{n_combined} slide** gabungan "
                    f"(1 Cover + {n_proposal_data} slide Proposal + "
                    f"{n_result_data} slide Result + 1 Penutup). "
                    f"Download tersedia di bawah."
                )
            except Exception as e:
                st.error(f"! Gagal: {e}")
                import traceback
                st.code(traceback.format_exc())

        # Tombol "Gabungkan" — hanya aktif kalau keduanya sudah di-generate terpisah
        both_ready = (
            "pptx_proposal_buf" in st.session_state
            and "pptx_result_buf" in st.session_state
        )
        if st.button(
            "! Gabungkan yang Sudah Di-generate",
            type="secondary",
            use_container_width=True,
            key="btn_gen_combined",
            disabled=not both_ready,
            help="Generate Proposal & Result terlebih dahulu (kolom di atas) sebelum menggabungkan."
                 if not both_ready else "Klik untuk menggabungkan kedua file yang sudah di-generate.",
        ):
            with st.spinner("Menggabungkan Proposal + Result..."):
                try:
                    combined_buf = merge_pptx(
                        st.session_state["pptx_proposal_buf"],
                        st.session_state["pptx_result_buf"],
                    )
                    st.session_state["pptx_combined_buf"] = combined_buf
                    st.session_state["pptx_combined_fname"] = (
                        f"Combined_{report_title.replace(' ', '_')}_"
                        f"{report_date.strftime('%Y%m%d')}.pptx"
                    )
                    n_proposal_data = df_proposal.shape[0]
                    n_result_data   = df_result.shape[0]
                    n_combined = 1 + n_proposal_data + n_result_data + 1
                    st.success(
                        f"✅ Berhasil digabungkan! "
                        f"Total **{n_combined} slide** gabungan "
                        f"(1 Cover + {n_proposal_data} slide Proposal + "
                        f"{n_result_data} slide Result + 1 Penutup)."
                    )
                except Exception as e:
                    st.error(f"! Gagal menggabungkan: {e}")
                    import traceback
                    st.code(traceback.format_exc())

        # Download gabungan — muncul setelah salah satu tombol di atas berhasil
        if "pptx_combined_buf" in st.session_state:
            st.download_button(
                label=" Download Gabungan .pptx",
                data=st.session_state["pptx_combined_buf"],
                file_name=st.session_state["pptx_combined_fname"],
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
                key="dl_combined",
            )

else:
    # Placeholder belum ada file
    st.markdown("""
    <div style="
        text-align: center;
        padding: 4rem 2rem;
        background: linear-gradient(135deg, #0d1117, #161b22);
        border: 1px dashed #CC0000;
        border-radius: 16px;
        color: #8b949e;
        margin-top: 1rem;
    ">
        <div style="font-size: 4rem; margin-bottom: 1rem;">📡</div>
        <h3 style="color: #e6edf3; margin-bottom: 0.5rem;">Belum ada file yang diupload</h3>
        <p>Upload file <strong>Excel (.xlsx / .xls)</strong> atau <strong>CSV</strong> di atas untuk memulai.</p>
        <p style="font-size:0.85rem; margin-top:1rem;">
             <em>Sheet <strong>Proposal</strong> (wajib): NO · SITE ID · SITE NAME · FINDING · CITY ·
            PURPOSE HEADER · SOW · PLAN ACTION · SUPPORT NEEDED · GOALS · INCREAMENT PAYLOAD AND REVENUE</em>
        </p>
        <p style="font-size:0.85rem; margin-top:0.5rem;">
             <em>Sheet <strong>Result</strong> (opsional): No · Background · Productivity Result</em>
        </p>
    </div>
    """, unsafe_allow_html=True)

    with st.expander(" Format Excel yang dibutuhkan"):
        tab_fmt_p, tab_fmt_r = st.tabs(["Sheet Proposal", "Sheet Result"])

        with tab_fmt_p:
            st.markdown("**Struktur kolom sheet Proposal:**")
            st.dataframe(pd.DataFrame([{
                "NO": "1",
                "SITE ID": "BDG212",
                "SITE NAME": "BDG212_SETRAMURNI-DMT",
                "FINDING": "Nearest Site BDG212_SETRAMURNI-DMT\nFrom BDG212 to Red Coverage",
                "CITY": "Bandung",
                "PURPOSE HEADER": "Install EM Under BDG212 SETRAMURNI",
                "SOW": "Purpose Easy macro 9 Meter -6.869313, 107.578078",
                "PLAN ACTION": "Propose Easy Macro 9 m -6.869313, 107.578078",
                "SUPPORT NEEDED": "Material AAU dual Band\nKabel Power + KWH",
                "GOALS": "Improvement payload & Revenue\nImprovement Red Coverage",
                "INCREAMENT PAYLOAD AND REVENUE": "Payload: 300 GB (Monthly)\nRevenue: 30 Jt (Monthly)",
            }]), use_container_width=True)

        with tab_fmt_r:
            st.markdown("**Struktur kolom sheet Result:**")
            st.dataframe(pd.DataFrame([{
                "No": "1",
                "Background": "ini isinya desc text background,\nabcdefghijklmnopqrstuvwxyz,\n12345678910...",
                "Productivity Result": "ini isinya desc text Productivity Result,\nabcdefghijklmnopqrstuvwxyz...",
            }]), use_container_width=True)
            st.info(
                " Header untuk PPT Result (Purpose Header & SOW) diambil otomatis dari "
                "sheet Proposal berdasarkan urutan baris yang sama."
            )
