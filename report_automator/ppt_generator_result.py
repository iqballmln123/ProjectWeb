"""
ppt_generator_result.py - Modul PPT untuk Sheet Result (v2)
=============================================================
Layout sesuai referensi user:

HEADER (maroon):
  ┌──────────────────────────────────────────────────────────────────────┐
  │ [○]  ┌─────────────────── white pill ───────────────────────────┐  Telkomsel │
  │      │ PURPOSE HEADER | SOW  (navy bold)                        │  (white)   │
  │      │ (Permanentizing Post Implementation Analysis) italic     │            │
  │      └──────────────────────────────────────────────────────────┘            │
  └──────────────────────────────────────────────────────────────────────┘

BODY (3 section dengan rounded container):

  ┌─────────────────┐  ┌──────────────────────────────┐  ┌──────────────────┐
  │   Background    │  │ Experience and Documentation │  │Productivity Result│
  │ ─────────────   │  │ Before                       │  │ RSRP & RSRQ After│
  │ [desc navy box] │  │ ┌─────────┬─────────┐        │  │ ┌──────────────┐ │
  │                 │  │ │  img 1  │  img 2  │        │  │ │   img A1     │ │
  │ Site Mapping    │  │ │         ├─────────┤        │  │ └──────────────┘ │
  │ ┌───────────┐   │  │ └─────────┤  img 3  │        │  │ ┌──────────────┐ │
  │ │  img SM   │   │  │           └─────────┘        │  │ │   img A2     │ │
  │ └───────────┘   │  │ After                        │  │ └──────────────┘ │
  │ RSRP(Before)    │  │ ┌─────────┬─────────┐        │  │ [desc navy box]  │
  │ ┌───────────┐   │  │ │  img 1  │  img 2  │        │  │                  │
  │ │ img RSRQ  │   │  │ │         ├─────────┤        │  │                  │
  │ └───────────┘   │  │ └─────────┤  img 3           │  │                  │
  └─────────────────┘  └──────────────────────────────┘  └──────────────────┘
"""

import io
from typing import List, Optional, Dict, Any
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from ppt_generator import (
    COLOR_WHITE, COLOR_BLACK, COLOR_DARK_NAVY, COLOR_GRAY_MID,
    COLOR_RED, COLOR_RED_DARK, COLOR_RED_LIGHT,
    SLIDE_W, SLIDE_H,
    COL_PURPOSE, COL_SOW, COL_SITE_NAME,
    _fill, _fill_border, _txb, _txt, _txt_multiline,
    _add_picture_fitted, _parse_bullets,
    _slide_1_cover, _slide_closing,
    _patch_theme_font,
)
from data_processor import COL_RESULT_BG, COL_RESULT_PR
from font_embedder import embed_poppins_into_pptx


# ─────────────────────────────────────────────────────────────────────────────
# WARNA KHUSUS RESULT
# ─────────────────────────────────────────────────────────────────────────────
COLOR_MAROON      = RGBColor(0x8B, 0x00, 0x00)   # Header maroon
COLOR_BODY_BG     = RGBColor(0xF0, 0xF2, 0xF5)   # Background body abu terang
COLOR_SEC_BORDER  = RGBColor(0x99, 0xAA, 0xBB)   # Border section container
COLOR_IMG_BG      = RGBColor(0xE0, 0xEA, 0xF6)   # Placeholder image bg (biru pucat)
COLOR_LABEL_TEXT  = RGBColor(0x1B, 0x2A, 0x4A)   # Label plain teks (navy)


# ─────────────────────────────────────────────────────────────────────────────
# KONSTANTA LAYOUT
# ─────────────────────────────────────────────────────────────────────────────
HDR_H    = 0.75   # tinggi header

BODY_TOP = HDR_H
BODY_BOT = 7.48   # hampir full height (sisa 0.02" footer tipis)
BODY_H   = BODY_BOT - BODY_TOP   # 6.73"

# Margin dan gap antar section
L_MARGIN = 0.08
SEC_GAP  = 0.055

# Lebar 3 kolom (total: 13.33" - 0.08*2 - 0.055*2 = 13.02")
COL1_W = 3.70   # Background
COL2_W = 5.62   # Experience & Documentation
COL3_W = 13.33 - L_MARGIN * 2 - SEC_GAP * 2 - COL1_W - COL2_W  # ≈ 3.72"

COL1_X = L_MARGIN                       # 0.08
COL2_X = COL1_X + COL1_W + SEC_GAP     # 3.835
COL3_X = COL2_X + COL2_W + SEC_GAP     # 9.51

# Ukuran section container
SEC_TOP = BODY_TOP + 0.04
SEC_H   = BODY_H - 0.06        # 6.67"

SEC_HDR_H   = 0.38   # tinggi header section (navy bar)
INNER_PX    = 0.07   # padding horizontal dalam section
INNER_PY    = 0.07   # padding vertikal setelah header section
ITEM_GAP    = 0.055  # gap antar elemen dalam section
LABEL_H     = 0.22   # tinggi plain label text


# ─────────────────────────────────────────────────────────────────────────────
# HELPER: Rounded Rectangle
# ─────────────────────────────────────────────────────────────────────────────
def _add_rounded_rect(slide, x: float, y: float, w: float, h: float,
                      fill_color: Optional[RGBColor] = None,
                      border_color: Optional[RGBColor] = None,
                      border_pt: float = 0.75,
                      corner_val: int = 5000,
                      transparent_fill: bool = False) -> Any:
    """
    Tambah rounded rectangle (shape type 5).
    corner_val: 0–50000, di mana 50000 = 50% radius.
    """
    shape = slide.shapes.add_shape(
        5,  # MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if transparent_fill:
        shape.fill.background()
    elif fill_color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if border_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    try:
        shape.adjustments[0] = corner_val
    except Exception:
        pass
    return shape


# ─────────────────────────────────────────────────────────────────────────────
# HELPER: Section container (rounded outer + navy header bar)
# ─────────────────────────────────────────────────────────────────────────────
def _draw_section_container(slide, x: float, y: float, w: float, h: float,
                            title: str) -> float:
    """
    Gambar section container:
    1. Outer rectangle biasa (transparent fill, navy border 2pt)
       — mulai dari header navy sampai bawah gambar
    2. Navy rectangle untuk header bar
    3. Teks judul section (putih, bold, centered)

    Returns: y coordinate awal content area (setelah header bar)
    """
    # 1. Outer rectangle (transparent, navy border 2pt)
    outer = slide.shapes.add_shape(
        1,                          # rectangle biasa
        Inches(x), Inches(y),
        Inches(w), Inches(h)
    )
    outer.fill.background()         # transparent
    outer.line.color.rgb = COLOR_DARK_NAVY
    outer.line.width = Pt(2.0)

    # 2. Navy header bar (rectangle, flush dengan outer border)
    hdr_rect = slide.shapes.add_shape(
        1,
        Inches(x), Inches(y),
        Inches(w), Inches(SEC_HDR_H)
    )
    hdr_rect.fill.solid()
    hdr_rect.fill.fore_color.rgb = COLOR_DARK_NAVY
    hdr_rect.line.fill.background()

    # 3. Section title text (white bold centered)
    txb = _txb(slide, x + 0.12, y + 0.07, w - 0.24, SEC_HDR_H - 0.08)
    _txt(txb.text_frame, title, 12, bold=True, color=COLOR_WHITE,
         align=PP_ALIGN.CENTER)

    return y + SEC_HDR_H + INNER_PY


# ─────────────────────────────────────────────────────────────────────────────
# HELPER: Plain label (small dark navy text, no background)
# ─────────────────────────────────────────────────────────────────────────────
def _draw_plain_label(slide, x: float, y: float, w: float, text: str,
                      size: float = 8.5):
    """Label teks kecil gelap tanpa background (seperti 'Site Mapping', 'Before')."""
    txb = _txb(slide, x, y, w, LABEL_H)
    _txt(txb.text_frame, text, size, bold=True, color=COLOR_LABEL_TEXT)


# ─────────────────────────────────────────────────────────────────────────────
# HELPER: Description box (navy fill, white centered text)
# ─────────────────────────────────────────────────────────────────────────────
def _draw_desc_box(slide, x: float, y: float, w: float, h: float, text: str):
    """Kotak transparent dengan teks navy, perataan kiri."""
    box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.background()   # transparent — tidak ada fill warna
    box.line.fill.background()

    lines = _parse_bullets(text)
    txb   = _txb(slide, x + 0.06, y + 0.06, w - 0.10, h - 0.08)
    tf    = txb.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT   # perataan kiri
        run = para.add_run()
        run.text = line.strip()
        run.font.size = Pt(10)
        run.font.color.rgb = COLOR_LABEL_TEXT   # navy gelap (terbaca di atas bg putih)
        run.font.name = "Poppins"


# ─────────────────────────────────────────────────────────────────────────────
# HELPER: Image slot (gambar atau placeholder dengan border)
# ─────────────────────────────────────────────────────────────────────────────
def _draw_img_slot(slide, img_bytes: Optional[bytes],
                   x: float, y: float, w: float, h: float,
                   placeholder_text: str):
    """Slot gambar: tampilkan gambar jika ada, atau placeholder berwarna."""
    # Border box background
    box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_IMG_BG
    box.line.color.rgb = COLOR_SEC_BORDER
    box.line.width = Pt(0.5)

    if img_bytes:
        _add_picture_fitted(slide, img_bytes, x, y, w, h, pad=0.04)
    else:
        txb = _txb(slide, x + 0.05, y + h / 2 - 0.15, w - 0.10, 0.30)
        _txt(txb.text_frame, placeholder_text, 7.5, italic=True,
             color=COLOR_GRAY_MID, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# HELPER: Grid 3 gambar  (kiri penuh | kanan split atas/bawah)
# ─────────────────────────────────────────────────────────────────────────────
def _draw_3img_grid(slide, imgs: List, x: float, y: float,
                    w: float, h: float, prefix: str, gap: float = 0.04):
    """
    Layout 3 gambar:
      ┌──────────┬──────────┐
      │          │  img 2   │
      │  img 1   ├──────────┤
      │          │  img 3   │
      └──────────┴──────────┘
    img1 → kiri penuh tinggi
    img2 → kanan atas (50% tinggi)
    img3 → kanan bawah (50% tinggi)
    """
    half_w = (w - gap) / 2
    half_h = (h - gap) / 2

    img1 = imgs[0] if len(imgs) > 0 else None
    img2 = imgs[1] if len(imgs) > 1 else None
    img3 = imgs[2] if len(imgs) > 2 else None

    _draw_img_slot(slide, img1, x, y, half_w, h, f"{prefix} 1")
    _draw_img_slot(slide, img2, x + half_w + gap, y, half_w, half_h, f"{prefix} 2")
    _draw_img_slot(slide, img3, x + half_w + gap, y + half_h + gap,
                   half_w, half_h, f"{prefix} 3")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE RESULT: Layout utama
# ─────────────────────────────────────────────────────────────────────────────
def _slide_result_site(
    prs,
    row_result: Dict[str, Any],
    row_proposal: Optional[Dict[str, Any]],
    images: Dict[str, List],
    slide_idx: int,
    total_slides: int,
):
    """
    Buat 1 slide Result sesuai layout referensi user.

    images dict keys:
      "site_mapping" : [bytes]       — 1 gambar
      "rsrp_before"  : [bytes]       — 1 gambar
      "before"       : [b1,b2,b3]   — 3 gambar
      "after"        : [b1,b2,b3]   — 3 gambar
      "rsrp_after"   : [b1,b2]      — 2 gambar
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ════════════════════════════════════════════════════════════════════════
    # LAYER 0: Backgrounds
    # ════════════════════════════════════════════════════════════════════════
    # Slide background putih
    bg_slide = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    _fill(bg_slide, COLOR_WHITE)

    # Body background abu terang
    bg_body = slide.shapes.add_shape(
        1, 0, Inches(BODY_TOP), SLIDE_W, Inches(BODY_H)
    )
    bg_body.fill.solid()
    bg_body.fill.fore_color.rgb = COLOR_BODY_BG
    bg_body.line.fill.background()

    # ════════════════════════════════════════════════════════════════════════
    # HEADER — Maroon background
    # ════════════════════════════════════════════════════════════════════════
    hdr_bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(HDR_H))
    hdr_bg.fill.solid()
    hdr_bg.fill.fore_color.rgb = COLOR_MAROON
    hdr_bg.line.fill.background()

    # ── Arrow navigasi kiri (NOTCHED_RIGHT_ARROW, putih di atas maroon) ───────
    arrow_nav = slide.shapes.add_shape(
        62,  # MSO_AUTO_SHAPE_TYPE.NOTCHED_RIGHT_ARROW
        Inches(0.10), Inches(0.115),
        Inches(0.55), Inches(0.52)
    )
    arrow_nav.fill.solid()
    arrow_nav.fill.fore_color.rgb = COLOR_WHITE
    arrow_nav.line.fill.background()

    # ── White pill (rounded rect) sebagai container judul ──────────────────
    pill = _add_rounded_rect(
        slide,
        x=0.73, y=0.065, w=10.82, h=0.62,
        fill_color=COLOR_WHITE,
        border_color=None,
        corner_val=8000,
    )

    # Data dari sheet Proposal untuk header
    if row_proposal:
        purpose = str(row_proposal.get(COL_PURPOSE, "")).strip()
        sow     = str(row_proposal.get(COL_SOW, "")).strip()
    else:
        purpose = sow = ""

    title_text    = f"{purpose}  |  {sow}" if purpose and sow else (purpose or "—")
    subtitle_text = "(Permanentizing Post Implementation Analysis)"

    # Baris 1: judul utama (navy bold) di dalam pill
    txb_title = _txb(slide, 0.88, 0.085, 10.40, 0.35)
    txb_title.text_frame.word_wrap = True
    _txt(txb_title.text_frame, title_text, 22, bold=True, color=COLOR_DARK_NAVY)

    # Baris 2: subtitle (navy italic)
    txb_sub = _txb(slide, 0.88, 0.41, 10.40, 0.22)
    _txt(txb_sub.text_frame, subtitle_text, 12, italic=True, color=COLOR_DARK_NAVY)

    # ── Telkomsel logo kanan (teks putih, vertikal center di header) ────────
    txb_tel = _txb(slide, 11.58, 0.17, 1.67, 0.42)
    _txt(txb_tel.text_frame, "Telkomsel", 18, bold=True, italic=True,
         color=COLOR_WHITE, align=PP_ALIGN.CENTER, font="Poppins")

    # ════════════════════════════════════════════════════════════════════════
    # SECTION 1 — Background
    # ════════════════════════════════════════════════════════════════════════
    content_y1 = _draw_section_container(
        slide, COL1_X, SEC_TOP, COL1_W, SEC_H, "Background"
    )

    # Batas bawah isi section 1 (dengan margin bawah)
    sec1_bot  = SEC_TOP + SEC_H - 0.05
    ix1       = COL1_X + INNER_PX
    iw1       = COL1_W - INNER_PX * 2
    iy1       = content_y1

    # ① Description text (navy box)
    DESC_H1 = 0.70
    bg_text = str(row_result.get(COL_RESULT_BG, "")).strip()
    _draw_desc_box(slide, ix1, iy1, iw1, DESC_H1, bg_text)
    iy1 += DESC_H1 + ITEM_GAP

    # ② Label "Site Mapping"
    _draw_plain_label(slide, ix1, iy1, iw1, "Site Mapping")
    iy1 += LABEL_H + 0.025

    # ③ Gambar Site Mapping
    SM_H   = 2.28
    sm_img = (images.get("site_mapping") or [None])[0]
    _draw_img_slot(slide, sm_img, ix1, iy1, iw1, SM_H, "[ Site Mapping ]")
    iy1 += SM_H + ITEM_GAP

    # ④ Label "RSRP & RSRQ (Before)"
    _draw_plain_label(slide, ix1, iy1, iw1, "RSRP & RSRQ (Before)")
    iy1 += LABEL_H + 0.025

    # ⑤ Gambar RSRP Before (isi sisa ruang)
    rsrq_b_h  = sec1_bot - iy1
    rsrq_b_img = (images.get("rsrp_before") or [None])[0]
    _draw_img_slot(slide, rsrq_b_img, ix1, iy1, iw1, rsrq_b_h,
                   "[ RSRP & RSRQ Before ]")

    # ════════════════════════════════════════════════════════════════════════
    # SECTION 2 — Experience and Documentation
    # ════════════════════════════════════════════════════════════════════════
    content_y2 = _draw_section_container(
        slide, COL2_X, SEC_TOP, COL2_W, SEC_H, "Experience and Documentation"
    )

    sec2_bot = SEC_TOP + SEC_H - 0.05
    ix2      = COL2_X + INNER_PX
    iw2      = COL2_W - INNER_PX * 2
    iy2      = content_y2

    # Hitung tinggi tiap grid (Before = 50%, After = 50% dari sisa tinggi)
    avail_h2    = sec2_bot - iy2
    each_grid_h = (avail_h2 - LABEL_H * 2 - ITEM_GAP * 3) / 2

    # ① Label "Before"
    _draw_plain_label(slide, ix2, iy2, iw2, "Before")
    iy2 += LABEL_H + 0.02

    # ② Grid Before (3 gambar)
    before_imgs = list(images.get("before", []))
    _draw_3img_grid(slide, before_imgs, ix2, iy2, iw2, each_grid_h, "Before")
    iy2 += each_grid_h + ITEM_GAP

    # ③ Label "After"
    _draw_plain_label(slide, ix2, iy2, iw2, "After")
    iy2 += LABEL_H + 0.02

    # ④ Grid After (3 gambar, isi sisa ruang)
    after_grid_h = sec2_bot - iy2
    after_imgs   = list(images.get("after", []))
    _draw_3img_grid(slide, after_imgs, ix2, iy2, iw2, after_grid_h, "After")

    # ════════════════════════════════════════════════════════════════════════
    # SECTION 3 — Productivity Result
    # ════════════════════════════════════════════════════════════════════════
    content_y3 = _draw_section_container(
        slide, COL3_X, SEC_TOP, COL3_W, SEC_H, "Productivity Result"
    )

    sec3_bot = SEC_TOP + SEC_H - 0.05
    ix3      = COL3_X + INNER_PX
    iw3      = COL3_W - INNER_PX * 2
    iy3      = content_y3

    # ① Label "RSRP & RSRQ (After)"
    _draw_plain_label(slide, ix3, iy3, iw3, "RSRP & RSRQ (After)")
    iy3 += LABEL_H + 0.02

    # Hitung tinggi 2 gambar + 1 desc box
    PROD_DESC_H = 1.38
    rsrq_total  = sec3_bot - iy3 - PROD_DESC_H - ITEM_GAP * 2
    each_rsrq_h = rsrq_total / 2

    # ② RSRQ After image 1
    rsrq_a_imgs = list(images.get("rsrp_after", []))
    img_a1      = rsrq_a_imgs[0] if len(rsrq_a_imgs) > 0 else None
    _draw_img_slot(slide, img_a1, ix3, iy3, iw3, each_rsrq_h,
                   "[ RSRP & RSRQ After 1 ]")
    iy3 += each_rsrq_h + ITEM_GAP

    # ③ RSRQ After image 2
    img_a2 = rsrq_a_imgs[1] if len(rsrq_a_imgs) > 1 else None
    _draw_img_slot(slide, img_a2, ix3, iy3, iw3, each_rsrq_h,
                   "[ RSRP & RSRQ After 2 ]")
    iy3 += each_rsrq_h + ITEM_GAP

    # ④ Description text Productivity Result (navy box, isi sisa)
    prod_desc_h = sec3_bot - iy3
    prod_text   = str(row_result.get(COL_RESULT_PR, "")).strip()
    _draw_desc_box(slide, ix3, iy3, iw3, prod_desc_h, prod_text)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE NUMBER (tipis, pojok kanan bawah)
    # ════════════════════════════════════════════════════════════════════════
    txb_pg = _txb(slide, 11.80, 7.38, 1.45, 0.12)
    _txt(txb_pg.text_frame,
         f"{slide_idx} / {total_slides}",
         6, color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────────────────────────────────────
# FUNGSI UTAMA: generate_result_pptx
# ─────────────────────────────────────────────────────────────────────────────
def generate_result_pptx(
    title: str,
    subtitle: str,
    author: str,
    report_date: str,
    df_proposal: pd.DataFrame,
    df_result: pd.DataFrame,
    site_result_images: Optional[Dict[str, Dict[str, List]]] = None,
) -> bytes:
    """
    Fungsi utama pembuatan PPTX Result Report Telkomsel.

    Sinkronisasi Proposal ↔ Result: index-based (baris ke-N di Result
    dipadankan baris ke-N di Proposal untuk mengambil PURPOSE HEADER & SOW).
    """
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    _patch_theme_font(prs, 'Poppins')   # ← ganti default theme font ke Poppins

    df_p = df_proposal.copy()
    df_p.columns    = [str(c).strip() for c in df_p.columns]
    proposal_rows   = df_p.to_dict(orient="records")

    df_r = df_result.copy()
    df_r.columns  = [str(c).strip() for c in df_r.columns]
    result_rows   = df_r.to_dict(orient="records")

    total_data   = len(result_rows)
    total_slides = 1 + total_data + 1   # cover + data + penutup

    # Slide 1: Cover
    _slide_1_cover(prs, title, subtitle, author, report_date)

    # Slide 2..N: satu per baris Result
    site_result_images = site_result_images or {}
    for i, row_res in enumerate(result_rows):
        row_prop = proposal_rows[i] if i < len(proposal_rows) else None
        images   = site_result_images.get(str(i), {})
        _slide_result_site(prs, row_res, row_prop, images, i + 2, total_slides)

    # Slide terakhir: Penutup
    _slide_closing(prs, author, report_date, title, total_slides)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    raw = buf.read()
    # Embed font Poppins ke dalam PPTX agar tampil benar di Canva / PowerPoint
    return embed_poppins_into_pptx(raw)
