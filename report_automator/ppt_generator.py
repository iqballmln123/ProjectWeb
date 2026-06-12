"""
ppt_generator.py - Modul pembuatan presentasi PowerPoint
==========================================================
Template khusus Telkomsel — Coverage Report Style.

Struktur slide:
  Slide 1       : Cover / Title
  Slide 2..N    : Satu slide per baris data (SITE ID → INCREMENT P&R)
  Slide terakhir: Penutup

Layout slide data (mirip gambar referensi):
  ┌──────────────────────────────────────────────────────────┐
  │  HEADER: Site Name | Purpose Header          [Telkomsel] │
  │  Sub-header merah: [City] Purpose ...                    │
  ├────────────────────┬─────────────────────────────────────┤
  │ ① FINDING          │  [Payload Site Surrounding Image]   │
  │   • bullet         │                                     │
  │ ② PLAN ACTION      ├─────────────────────────────────────┤
  │   teks             │  [Maps Preview & COVMO]             │
  │ ③ SUPPORT NEEDED   │                                     │
  │   • bullet         ├──────────────────┬──────────────────┤
  │ ④ GOALS            │  [Support Img 1] │ [Support Img 2]  │
  │   • bullet         │                  │                  │
  │ ⑤ INCREMENT P&R    └──────────────────┴──────────────────┘
  │   • bullet         (jika support img hanya 1: full lebar)
  └────────────────────┘
"""

import io
import math
from typing import List, Optional, Dict, Any
from lxml import etree
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image as PILImage
from font_embedder import embed_poppins_into_pptx


# ─────────────────────────────────────────────────────────────────────────────
# KONSTANTA: Palet warna Telkomsel (sesuai gambar referensi)
# ─────────────────────────────────────────────────────────────────────────────
COLOR_RED        = RGBColor(0xCC, 0x00, 0x00)   # #CC0000 — Merah Telkomsel
COLOR_RED_DARK   = RGBColor(0x99, 0x00, 0x00)   # #990000 — Merah gelap
COLOR_RED_LIGHT  = RGBColor(0xFF, 0x33, 0x33)   # #FF3333 — Merah terang
COLOR_RED_MAROON = RGBColor(0xCC, 0x00, 0x33)   # #CC0033 — Merah maroon (sub-header)
COLOR_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BLACK      = RGBColor(0x00, 0x00, 0x00)
COLOR_DARK_NAVY  = RGBColor(0x1B, 0x2A, 0x4A)   # #1B2A4A — Blue navy gelap (header)
COLOR_GRAY_BG    = RGBColor(0xF5, 0xF5, 0xF5)   # Background abu terang
COLOR_GRAY_MID   = RGBColor(0xAA, 0xAA, 0xAA)
COLOR_GOLD       = RGBColor(0xFF, 0xC4, 0x00)   # Aksen emas
COLOR_ORANGE_SUB = RGBColor(0xFF, 0x66, 0x00)   # Sub-header oranye/kuning
COLOR_SECTION_BG = RGBColor(0xCC, 0x00, 0x00)   # Background label seksi (merah)
COLOR_CONTENT_BG = RGBColor(0xFF, 0xFF, 0xFF)   # Background isi konten

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# Pemetaan kolom Excel → field internal
COL_NO           = "NO"
COL_SITE_ID      = "SITE ID"
COL_SITE_NAME    = "SITE NAME"
COL_FINDING      = "FINDING"
COL_CITY         = "CITY"
COL_PURPOSE      = "PURPOSE HEADER"
COL_SOW          = "SOW"
COL_PLAN_ACTION  = "PLAN ACTION"
COL_SUPPORT      = "SUPPORT NEEDED"
COL_GOALS        = "GOALS"
COL_INCREMENT    = "INCREAMENT PAYLOAD AND REVENUE"


# ─────────────────────────────────────────────────────────────────────────────
# HELPERS dasar
# ─────────────────────────────────────────────────────────────────────────────
def _fill(shape, color: RGBColor):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color
    shape.line.fill.background()


def _fill_border(shape, fill_color: RGBColor, border_color: RGBColor, border_pt: float = 0.75):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(border_pt)


def _rect(slide, left, top, width, height, color: RGBColor):
    """Tambah rectangle dengan solid fill. Semua ukuran dalam inches."""
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    _fill(shape, color)
    return shape


def _txb(slide, left, top, width, height):
    """Tambah textbox. Semua ukuran dalam inches."""
    return slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )


def _txt(
    tf,
    text: str,
    size: float,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = COLOR_BLACK,
    align=PP_ALIGN.LEFT,
    font: str = "Poppins",
):
    """Set teks di text frame (clear dulu)."""
    tf.clear()
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font


def _txt_multiline(
    tf,
    lines: List[str],
    size: float,
    bold: bool = False,
    color: RGBColor = COLOR_BLACK,
    bullet: bool = False,
    font: str = "Poppins",
):
    """Set multi-baris teks di text frame."""
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        if bullet:
            para.space_before = Pt(1)
        run = para.add_run()
        prefix = "• " if bullet and line.strip() else ""
        run.text = prefix + line.strip()
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font


def _patch_theme_font(prs, font_name: str = "Poppins"):
    """
    Patch theme XML presentasi agar major/minor font → Poppins.
    Ini mencegah fallback ke Times New Roman / Calibri dari default template.
    """
    from lxml import etree as _etree
    NSMAP = "http://schemas.openxmlformats.org/drawingml/2006/main"
    for theme_part in prs.part.part_related_by(
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
    ).blob and [] or []:
        pass  # tidak dipakai

    try:
        # Akses theme part langsung via slide_master
        slide_master = prs.slide_masters[0]
        theme_part   = slide_master.part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        )
        root = _etree.fromstring(theme_part.blob)
        ns   = {"a": NSMAP}

        # Cari elemen majorFont dan minorFont dalam themeElements → fmtScheme → ...
        for tag in ("majorFont", "minorFont"):
            nodes = root.findall(f".//a:{tag}", ns)
            for node in nodes:
                # Set latin typeface
                latin = node.find("a:latin", ns)
                if latin is None:
                    latin = _etree.SubElement(
                        node, f"{{{NSMAP}}}latin"
                    )
                latin.set("typeface", font_name)

        # Simpan kembali blob yang sudah dipatch
        theme_part._blob = _etree.tostring(root, xml_declaration=True,
                                           encoding="UTF-8", standalone=True)
    except Exception:
        # Fallback jika struktur theme berbeda — abaikan saja
        pass


def _add_picture_fitted(slide, img_bytes: bytes,
                         left_in: float, top_in: float,
                         w_in: float, h_in: float,
                         pad: float = 0.05):
    """
    Tambah gambar ke slide dengan aspect-ratio preserved dan center-aligned
    dalam bounding box yang diberikan. Semua dalam inches.
    """
    avail_w = int(Inches(w_in - 2 * pad))
    avail_h = int(Inches(h_in - 2 * pad))

    try:
        img = PILImage.open(io.BytesIO(img_bytes))
        iw, ih = img.size
        asp = iw / ih
        if int(avail_w / asp) <= avail_h:
            fw, fh = avail_w, int(avail_w / asp)
        else:
            fh, fw = avail_h, int(avail_h * asp)
        ox = (avail_w - fw) // 2
        oy = (avail_h - fh) // 2
        pl = int(Inches(left_in + pad)) + ox
        pt = int(Inches(top_in + pad)) + oy
    except Exception:
        fw, fh = avail_w, avail_h
        pl = int(Inches(left_in + pad))
        pt = int(Inches(top_in + pad))

    slide.shapes.add_picture(io.BytesIO(img_bytes), pl, pt, fw, fh)


def _parse_bullets(text: str) -> List[str]:
    """
    Pecah teks berdasarkan newline atau koma/titik-koma menjadi list bullet.
    """
    if not text or str(text).strip() in ("", "nan", "None"):
        return ["—"]
    text = str(text).strip()
    # Pisahkan berdasarkan newline atau \n literal
    lines = [l.strip() for l in text.replace("\\n", "\n").splitlines()]
    # Hapus baris kosong
    lines = [l for l in lines if l]
    return lines if lines else ["—"]


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1: Cover
# ─────────────────────────────────────────────────────────────────────────────
def _slide_1_cover(prs, title: str, subtitle: str, author: str, report_date: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background: kiri merah, kanan gelap
    left_panel = slide.shapes.add_shape(1, 0, 0, Inches(7.3), SLIDE_H)
    _fill(left_panel, COLOR_RED_DARK)

    right_panel = slide.shapes.add_shape(1, Inches(7.3), 0, Inches(6.03), SLIDE_H)
    _fill(right_panel, COLOR_DARK_NAVY)

    # Garis vertikal emas
    div = slide.shapes.add_shape(1, Inches(7.27), 0, Inches(0.07), SLIDE_H)
    _fill(div, COLOR_GOLD)

    # Strip atas & bawah
    for y in (0, Inches(7.38)):
        s = slide.shapes.add_shape(1, 0, y, SLIDE_W, Inches(0.12))
        _fill(s, COLOR_RED_LIGHT)

    # Logo teks kiri
    txb = _txb(slide, 0.5, 0.3, 5.5, 0.75)
    _txt(txb.text_frame, "TELKOMSEL", 32, bold=True, color=COLOR_WHITE, font="Poppins")

    txb_tag = _txb(slide, 0.5, 0.98, 6.0, 0.35)
    _txt(txb_tag.text_frame, "The Biggest & Most Innovative Digital Telco of Indonesia",
         9, italic=True, color=RGBColor(0xFF, 0xCC, 0xCC))

    # Garis emas horizontal
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.05))
    _fill(line, COLOR_GOLD)

    # Judul
    txb_title = _txb(slide, 0.5, 1.7, 6.5, 2.8)
    txb_title.text_frame.word_wrap = True
    _txt(txb_title.text_frame, title, 30, bold=True, color=COLOR_WHITE)

    # Sub-judul
    txb_sub = _txb(slide, 0.5, 4.45, 6.5, 0.7)
    txb_sub.text_frame.word_wrap = True
    _txt(txb_sub.text_frame, subtitle, 14, italic=True, color=RGBColor(0xFF, 0xCC, 0xCC))

    # Panel kanan: info cards
    badge = slide.shapes.add_shape(1, Inches(7.9), Inches(1.0), Inches(5.0), Inches(0.45))
    _fill(badge, COLOR_RED)
    txb_b = _txb(slide, 7.98, 1.07, 4.8, 0.32)
    _txt(txb_b.text_frame, "▶  LAPORAN RESMI — INTERNAL USE ONLY", 9, bold=True, color=COLOR_WHITE)

    info_items = [
        ("Tanggal Laporan", report_date),
        ("Disusun Oleh", author),
        ("Unit Kerja", subtitle),
        ("Jenis Laporan", "Coverage Activity Report"),
    ]
    y_info = 1.65
    for label, val in info_items:
        card = slide.shapes.add_shape(1, Inches(7.55), Inches(y_info), Inches(5.4), Inches(0.72))
        _fill(card, RGBColor(0x25, 0x25, 0x45))
        card.line.color.rgb = COLOR_RED

        txb_l = _txb(slide, 7.66, y_info + 0.07, 5.1, 0.25)
        _txt(txb_l.text_frame, label, 8, color=RGBColor(0xAA, 0xAA, 0xCC))

        txb_v = _txb(slide, 7.66, y_info + 0.33, 5.1, 0.33)
        _txt(txb_v.text_frame, val, 12, bold=True, color=COLOR_WHITE)

        y_info += 0.84

    txb_cl = _txb(slide, 7.5, 6.82, 5.5, 0.32)
    _txt(txb_cl.text_frame, "RAHASIA — TIDAK UNTUK DISEBARLUASKAN",
         8, bold=True, color=COLOR_RED_LIGHT, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE DATA: Layout 5-seksi kiri + area gambar kanan
# ─────────────────────────────────────────────────────────────────────────────
def _draw_section_label(slide, num: int, label: str,
                        left: float, top: float, width: float, height: float = 0.26):
    """
    Gambar label seksi bernomor (kotak merah + teks).
    Mirip bagian '① FINDING' di gambar referensi.
    """
    # Badge nomor (kotak merah kecil)
    badge_w = 0.26
    badge = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(badge_w), Inches(height)
    )
    _fill(badge, COLOR_RED)

    txb_num = _txb(slide, left + 0.01, top, badge_w - 0.02, height)
    _txt(txb_num.text_frame, str(num), 9, bold=True,
         color=COLOR_WHITE, align=PP_ALIGN.CENTER)

    # Label teks merah bold
    label_bg = slide.shapes.add_shape(
        1, Inches(left + badge_w), Inches(top),
        Inches(width - badge_w), Inches(height)
    )
    _fill(label_bg, COLOR_RED)

    txb_lbl = _txb(slide, left + badge_w + 0.04, top + 0.02,
                   width - badge_w - 0.06, height - 0.04)
    _txt(txb_lbl.text_frame, label, 9, bold=True, color=COLOR_WHITE)


def _draw_content_area(slide, lines: List[str],
                       left: float, top: float,
                       width: float, height: float,
                       font_size: float = 18,
                       bullet: bool = True):
    """
    Gambar area konten teks (putih/abu muda) dengan bullet points.
    """
    bg = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    _fill(bg, COLOR_CONTENT_BG)

    txb = _txb(slide, left + 0.06, top + 0.04, width - 0.1, height - 0.08)
    txb.text_frame.word_wrap = True
    _txt_multiline(txb.text_frame, lines, font_size,
                   color=COLOR_BLACK, bullet=bullet)



def _draw_right_label(slide, num_str: str, label: str,
                      x: float, y: float, w: float, h: float = 0.24):
    """
    Label area kanan (⑥ PAYLOAD, ⑦ MAPS) dengan badge merah kecil + teks putih bold.
    """
    lbl_bg = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    _fill(lbl_bg, COLOR_DARK_NAVY)

    # Badge merah kecil
    badge_w = 0.20
    badge = slide.shapes.add_shape(
        1, Inches(x + 0.06), Inches(y + 0.02), Inches(badge_w), Inches(h - 0.04)
    )
    _fill(badge, COLOR_RED)
    txb_n = _txb(slide, x + 0.06, y + 0.02, badge_w, h - 0.04)
    _txt(txb_n.text_frame, num_str, 7.5, bold=True,
         color=COLOR_WHITE, align=PP_ALIGN.CENTER)

    # Label teks
    txb_lbl = _txb(slide, x + 0.06 + badge_w + 0.06, y + 0.03,
                   w - 0.06 - badge_w - 0.12, h - 0.06)
    _txt(txb_lbl.text_frame, label, 9, bold=True, color=COLOR_WHITE)


def _add_native_pic_placeholder(
    slide,
    idx: int,
    left_in: float,
    top_in: float,
    w_in: float,
    h_in: float,
    label: str = "Double-click to insert image",
):
    """
    Inject native PowerPoint picture placeholder ke dalam slide.

    User bisa langsung double-click frame ini di PowerPoint untuk
    membuka file browser dan insert gambar — tanpa perlu edit XML manual.

    Args:
        slide   : slide object python-pptx
        idx     : unique integer per placeholder dalam satu slide (mulai 10)
        left_in : posisi kiri (inches)
        top_in  : posisi atas (inches)
        w_in    : lebar (inches)
        h_in    : tinggi (inches)
        label   : nama placeholder (muncul di panel Properties PowerPoint)
    """
    import xml.sax.saxutils as saxutils

    spTree = slide.shapes._spTree

    left        = int(Inches(left_in))
    top         = int(Inches(top_in))
    width       = int(Inches(w_in))
    height      = int(Inches(h_in))
    safe_label  = saxutils.escape(label)  # escape & < > untuk XML-safe

    # XML element untuk native picture placeholder (type="pic")
    # PowerPoint mengenali ini sebagai slot gambar yang bisa di-double-click
    sp_xml = (
        f'<p:sp '
        f'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        f'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        f'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        f'<p:nvSpPr>'
        f'<p:cNvPr id="{idx + 200}" name="{safe_label}"/>'
        f'<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
        f'<p:nvPr><p:ph type="pic" idx="{idx}"/></p:nvPr>'
        f'</p:nvSpPr>'
        f'<p:spPr>'
        f'<a:xfrm><a:off x="{left}" y="{top}"/><a:ext cx="{width}" cy="{height}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f'<a:solidFill><a:srgbClr val="F0F0F0"/></a:solidFill>'
        f'<a:ln w="12700"><a:solidFill><a:srgbClr val="CC0000"/></a:solidFill>'
        f'<a:prstDash val="dashDot"/></a:ln>'
        f'</p:spPr>'
        f'<p:txBody>'
        f'<a:bodyPr anchor="ctr"/>'
        f'<a:lstStyle/>'
        f'<a:p><a:pPr algn="ctr"/>'
        f'<a:r><a:rPr lang="id-ID" sz="1000" i="1" dirty="0"/>'
        f'<a:t>\U0001f5bc\ufe0f  {safe_label}</a:t></a:r></a:p>'
        f'</p:txBody>'
        f'</p:sp>'
    )

    sp_element = etree.fromstring(sp_xml)
    spTree.append(sp_element)


def _slide_data_site(
    prs,
    row: Dict[str, Any],
    images: List[Dict],  # list of {"bytes": bytes, "label": str}
    slide_idx: int,
    total_slides: int,
):
    """
    Membuat 1 slide data untuk 1 baris (1 site) dari Excel.

    Layout FINAL (sesuai gambar referensi):
    ──────────────────────────────────────────────────────────────────────
    Y=0.00 │ HEADER BAR: "PURPOSE HEADER"  bold hitam     [Telkomsel®]  │
    Y=0.52 │ SUB-HEADER: "[CITY] SOW..."   kuning di atas navy          │
    Y=0.82 │ Panel KIRI (4.55")  │ Sub-kiri Payload+Maps │ Sub-kanan   │
           │ ① FINDING           │ ⑥ label               │             │
           │   • bullet          │   [payload chart]      │  [support   │
           │ ② PLAN ACTION       │                        │   img 1]    │
           │   teks              ├───────────────────────│             │
           │ ③ SUPPORT NEEDED    │ ⑦ label               │─────────────│
           │   • bullet          │   [maps/COVMO]         │  [supp 2]   │
           │ ④ GOALS             │                        │  [supp 3]   │
           │   - bullet          │                        │             │
           │ ⑤ INCREMENT P&R     │                        │             │
           │   • bullet          │                        │             │
    Y=7.28 │ FOOTER bar (navy)                                           │
    ──────────────────────────────────────────────────────────────────────

    images[0] = Payload Site Surrounding chart   → sub-kiri atas
    images[1] = Maps Preview & COVMO             → sub-kiri bawah
    images[2] = support image 1                  → sub-kanan atas  (row A)
    images[3] = support image 2                  → sub-kanan bawah kiri (row B)
    images[4] = support image 3                  → sub-kanan bawah kanan (row B)

    Jika support img hanya 1 → isi penuh kolom sub-kanan
    Jika support img 2       → img1 top (row A tinggi), img2 bottom (row B tinggi)
    Jika support img 3       → img1 top, img2+img3 grid 2-kolom di bottom
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ── Background putih ────────────────────────────────────────────────────
    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    _fill(bg, COLOR_WHITE)

    # ── KONSTANTA LAYOUT ────────────────────────────────────────────────────
    HEADER_H  = 0.82   # tinggi header gabungan (navy bar sepenuhnya)
    FOOTER_H  = 0.22   # tinggi footer
    CONTENT_T = HEADER_H          # 0.82 — awal area konten
    CONTENT_B = 7.5 - FOOTER_H   # 7.28 — akhir area konten
    CONTENT_H = CONTENT_B - CONTENT_T  # 6.46"

    LEFT_W    = 4.55   # lebar panel kiri (5 seksi)
    RIGHT_X   = LEFT_W
    RIGHT_W   = 13.33 - LEFT_W  # 8.78" — total panel kanan

    # Sub-kolom dalam panel kanan
    LEFT_SUB_W  = 5.30  # lebar sub-kiri (payload + maps)
    RIGHT_SUB_X = RIGHT_X + LEFT_SUB_W   # 9.85
    RIGHT_SUB_W = RIGHT_W - LEFT_SUB_W   # 3.48"

    # Pembagian baris (row A = ~45%, row B = ~55%)
    ROW_A_H = CONTENT_H * 0.45   # ≈2.907"
    ROW_B_H = CONTENT_H - ROW_A_H
    ROW_A_T = CONTENT_T
    ROW_B_T = CONTENT_T + ROW_A_H

    LABEL_H = 0.24  # tinggi label seksi kanan (⑥ dan ⑦)

    # ── HEADER BAR — satu bar navy biru gelap ─────────────────────────────
    # Background navy penuh
    hdr_bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(HEADER_H))
    _fill(hdr_bg, COLOR_DARK_NAVY)

    # Data teks
    purpose      = str(row.get(COL_PURPOSE, "")).strip()
    site_name    = str(row.get(COL_SITE_NAME, "")).strip()
    city         = str(row.get(COL_CITY, "")).strip()
    sow          = str(row.get(COL_SOW, "")).strip()
    plan_action  = str(row.get(COL_PLAN_ACTION, "")).strip().split("\n")[0]  # baris pertama saja

    # Baris 1: "PURPOSE HEADER | SOW" — putih, bold
    if purpose and sow:
        header_line1 = f"{purpose}  |  {sow}"
    elif purpose:
        header_line1 = purpose
    else:
        header_line1 = site_name

    # Baris 2: "[City] PLAN ACTION" — merah maroon, italic
    header_line2 = f"[{city}]  {plan_action}" if plan_action else f"[{city}]"

    # Textbox baris 1 (putih bold)
    txb_line1 = _txb(slide, 0.18, 0.07, 10.6, 0.44)
    txb_line1.text_frame.word_wrap = True
    _txt(txb_line1.text_frame, header_line1, 18, bold=True, color=COLOR_WHITE)

    # Textbox baris 2 (putih bold)
    txb_line2 = _txb(slide, 0.18, 0.50, 10.6, 0.28)
    txb_line2.text_frame.word_wrap = False
    _txt(txb_line2.text_frame, header_line2, 12, bold=True, color=COLOR_WHITE)

    # Logo "Telkomsel" pojok kanan atas (merah, italic, Poppins)
    txb_logo = _txb(slide, 10.85, 0.08, 2.35, 0.42)
    _txt(txb_logo.text_frame, "Telkomsel", 20, bold=True, italic=True,
         color=COLOR_RED, align=PP_ALIGN.RIGHT, font="Poppins")

    # ── GARIS VERTIKAL PEMISAH KIRI-KANAN ───────────────────────────────
    vline_main = slide.shapes.add_shape(
        1, Inches(LEFT_W - 0.02), Inches(CONTENT_T),
        Inches(0.03), Inches(CONTENT_H)
    )
    _fill(vline_main, COLOR_RED)

    # ── GARIS VERTIKAL PEMISAH SUB-KIRI & SUB-KANAN ─────────────────────
    vline_sub = slide.shapes.add_shape(
        1, Inches(RIGHT_SUB_X - 0.02), Inches(CONTENT_T),
        Inches(0.03), Inches(CONTENT_H)
    )
    _fill(vline_sub, COLOR_RED)

    # ── GARIS HORIZONTAL PEMISAH ROW A & ROW B (hanya di sub-kiri) ──────
    hsep = slide.shapes.add_shape(
        1, Inches(RIGHT_X), Inches(ROW_B_T - 0.015),
        Inches(LEFT_SUB_W), Inches(0.03)
    )
    _fill(hsep, COLOR_RED)

    # ────────────────────────────────────────────────────────────────────
    # PANEL KIRI: 5 seksi
    # ────────────────────────────────────────────────────────────────────
    sections = [
        {
            "num": 1, "label": "FINDING",
            "content": _parse_bullets(row.get(COL_FINDING, "")),
            "bullet": True,
        },
        {
            "num": 2, "label": "PLAN ACTION",
            "content": _parse_bullets(row.get(COL_PLAN_ACTION, "")),
            "bullet": False,
        },
        {
            "num": 3, "label": "SUPPORT NEEDED",
            "content": _parse_bullets(row.get(COL_SUPPORT, "")),
            "bullet": True,
        },
        {
            "num": 4, "label": "GOALS",
            "content": _parse_bullets(row.get(COL_GOALS, "")),
            "bullet": True,
        },
        {
            "num": 5, "label": "Increment Payload & Revenue",
            "content": _parse_bullets(row.get(COL_INCREMENT, "")),
            "bullet": True,
        },
    ]

    SEC_LABEL_H = 0.25   # tinggi label seksi kiri
    SEC_PAD     = 0.035  # jarak antar elemen kiri

    # Hitung distribusi tinggi konten per seksi
    total_label_h  = len(sections) * SEC_LABEL_H
    total_pad_h    = len(sections) * SEC_PAD * 2
    avail_h        = CONTENT_H - total_label_h - total_pad_h
    total_lines    = max(sum(len(s["content"]) for s in sections), 1)

    # Bobot minimum per seksi agar tidak terlalu kecil
    MIN_SEC_H = 0.28

    y_cur = CONTENT_T
    for sec in sections:
        n_lines = len(sec["content"])
        raw_h   = (n_lines / total_lines) * avail_h
        sec_h   = max(raw_h, MIN_SEC_H)

        # ── Label seksi ─────────────────────────────────────
        lbl_y  = y_cur + SEC_PAD
        badge_w = 0.25
        badge_h = SEC_LABEL_H

        # Badge nomor merah
        badge = slide.shapes.add_shape(
            1, Inches(0.05), Inches(lbl_y), Inches(badge_w), Inches(badge_h)
        )
        _fill(badge, COLOR_RED)
        txb_n = _txb(slide, 0.05, lbl_y, badge_w, badge_h)
        _txt(txb_n.text_frame, str(sec["num"]), 8, bold=True,
             color=COLOR_WHITE, align=PP_ALIGN.CENTER)

        # Label teks (merah bg)
        lbl_bg = slide.shapes.add_shape(
            1, Inches(0.05 + badge_w), Inches(lbl_y),
            Inches(LEFT_W - 0.10 - badge_w), Inches(badge_h)
        )
        _fill(lbl_bg, COLOR_RED)
        txb_lbl = _txb(slide, 0.05 + badge_w + 0.05, lbl_y + 0.02,
                       LEFT_W - 0.15 - badge_w, badge_h - 0.04)
        _txt(txb_lbl.text_frame, sec["label"], 8.5, bold=True, color=COLOR_WHITE)

        y_cur += SEC_PAD + badge_h

        # ── Konten seksi ────────────────────────────────────
        cont_bg = slide.shapes.add_shape(
            1, Inches(0.05), Inches(y_cur),
            Inches(LEFT_W - 0.10), Inches(sec_h)
        )
        _fill(cont_bg, COLOR_WHITE)

        txb_c = _txb(slide, 0.15, y_cur + 0.03, LEFT_W - 0.22, sec_h - 0.05)
        txb_c.text_frame.word_wrap = True
        _txt_multiline(txb_c.text_frame, sec["content"], 12,
                       color=COLOR_BLACK, bullet=sec["bullet"])

        y_cur += sec_h + SEC_PAD

    # ────────────────────────────────────────────────────────────────────
    # PANEL KANAN — SUB-KIRI: Payload (atas) + Maps (bawah)
    # ────────────────────────────────────────────────────────────────────

    # Background sub-kiri
    sub_left_bg = slide.shapes.add_shape(
        1, Inches(RIGHT_X), Inches(CONTENT_T),
        Inches(LEFT_SUB_W), Inches(CONTENT_H)
    )
    _fill(sub_left_bg, COLOR_GRAY_BG)

    # ── LABEL ⑥ PAYLOAD Site Surrounding ──────────────────────────────
    _draw_right_label(slide, "6", "PAYLOAD site Surrounding",
                      x=RIGHT_X, y=ROW_A_T, w=LEFT_SUB_W, h=LABEL_H)

    # Gambar payload
    img_payload_top = ROW_A_T + LABEL_H
    img_payload_h   = ROW_A_H - LABEL_H

    if len(images) >= 1 and images[0].get("bytes"):
        _add_picture_fitted(
            slide, images[0]["bytes"],
            left_in=RIGHT_X, top_in=img_payload_top,
            w_in=LEFT_SUB_W, h_in=img_payload_h, pad=0.05
        )
    else:
        _add_native_pic_placeholder(
            slide, idx=10,
            left_in=RIGHT_X, top_in=img_payload_top,
            w_in=LEFT_SUB_W, h_in=img_payload_h,
            label="Payload Site Surrounding",
        )

    # ── LABEL ⑦ MAPS Preview & COVMO ──────────────────────────────────
    _draw_right_label(slide, "7", "MAPS Preview & COVMO",
                      x=RIGHT_X, y=ROW_B_T, w=LEFT_SUB_W, h=LABEL_H)

    # Gambar maps
    img_maps_top = ROW_B_T + LABEL_H
    img_maps_h   = ROW_B_H - LABEL_H

    if len(images) >= 2 and images[1].get("bytes"):
        _add_picture_fitted(
            slide, images[1]["bytes"],
            left_in=RIGHT_X, top_in=img_maps_top,
            w_in=LEFT_SUB_W, h_in=img_maps_h, pad=0.05
        )
    else:
        _add_native_pic_placeholder(
            slide, idx=11,
            left_in=RIGHT_X, top_in=img_maps_top,
            w_in=LEFT_SUB_W, h_in=img_maps_h,
            label="Maps Preview & COVMO",
        )

    # ────────────────────────────────────────────────────────────────────
    # PANEL KANAN — SUB-KANAN: Support images
    # ────────────────────────────────────────────────────────────────────
    # Background sub-kanan
    sub_right_bg = slide.shapes.add_shape(
        1, Inches(RIGHT_SUB_X), Inches(CONTENT_T),
        Inches(RIGHT_SUB_W), Inches(CONTENT_H)
    )
    _fill(sub_right_bg, COLOR_GRAY_BG)

    support_imgs = [img for img in images[2:5] if img.get("bytes")]
    n_sup = len(support_imgs)

    if n_sup == 0:
        # Native picture placeholder penuh — user bisa double-click di PPT
        _add_native_pic_placeholder(
            slide, idx=12,
            left_in=RIGHT_SUB_X, top_in=CONTENT_T,
            w_in=RIGHT_SUB_W, h_in=CONTENT_H,
            label="Support Image",
        )

    elif n_sup == 1:
        # 1 gambar → isi penuh kolom sub-kanan
        _add_picture_fitted(
            slide, support_imgs[0]["bytes"],
            left_in=RIGHT_SUB_X, top_in=CONTENT_T,
            w_in=RIGHT_SUB_W, h_in=CONTENT_H, pad=0.05
        )

    elif n_sup == 2:
        # 2 gambar → img1 atas (sejajar Row A), img2 bawah (sejajar Row B)
        # Garis horizontal pemisah
        hsep_r = slide.shapes.add_shape(
            1, Inches(RIGHT_SUB_X), Inches(ROW_B_T - 0.015),
            Inches(RIGHT_SUB_W), Inches(0.03)
        )
        _fill(hsep_r, COLOR_RED)

        _add_picture_fitted(
            slide, support_imgs[0]["bytes"],
            left_in=RIGHT_SUB_X, top_in=ROW_A_T,
            w_in=RIGHT_SUB_W, h_in=ROW_A_H, pad=0.05
        )
        _add_picture_fitted(
            slide, support_imgs[1]["bytes"],
            left_in=RIGHT_SUB_X, top_in=ROW_B_T,
            w_in=RIGHT_SUB_W, h_in=ROW_B_H, pad=0.05
        )

    else:  # n_sup == 3
        # 3 gambar:
        #   img1 → atas, full lebar sub-kanan (sejajar Row A)
        #   img2 + img3 → bawah, 2 kolom (sejajar Row B)
        half_w = RIGHT_SUB_W / 2

        # Garis horizontal
        hsep_r = slide.shapes.add_shape(
            1, Inches(RIGHT_SUB_X), Inches(ROW_B_T - 0.015),
            Inches(RIGHT_SUB_W), Inches(0.03)
        )
        _fill(hsep_r, COLOR_RED)

        # Garis vertikal antara img2 dan img3
        vsep_r = slide.shapes.add_shape(
            1, Inches(RIGHT_SUB_X + half_w - 0.015), Inches(ROW_B_T),
            Inches(0.03), Inches(ROW_B_H)
        )
        _fill(vsep_r, COLOR_RED)

        _add_picture_fitted(
            slide, support_imgs[0]["bytes"],
            left_in=RIGHT_SUB_X, top_in=ROW_A_T,
            w_in=RIGHT_SUB_W, h_in=ROW_A_H, pad=0.05
        )
        _add_picture_fitted(
            slide, support_imgs[1]["bytes"],
            left_in=RIGHT_SUB_X, top_in=ROW_B_T,
            w_in=half_w, h_in=ROW_B_H, pad=0.04
        )
        _add_picture_fitted(
            slide, support_imgs[2]["bytes"],
            left_in=RIGHT_SUB_X + half_w, top_in=ROW_B_T,
            w_in=half_w, h_in=ROW_B_H, pad=0.04
        )

    # ── FOOTER BAR ──────────────────────────────────────────────────────
    footer_top = 7.5 - FOOTER_H
    footer_bg = slide.shapes.add_shape(
        1, 0, Inches(footer_top), SLIDE_W, Inches(FOOTER_H)
    )
    _fill(footer_bg, COLOR_DARK_NAVY)

    txb_ft_l = _txb(slide, 0.2, footer_top + 0.03, 7.0, FOOTER_H - 0.04)
    _txt(txb_ft_l.text_frame, "PT. Telkomsel — Connecting Indonesia",
         7, color=RGBColor(0xAA, 0xAA, 0xBB))

    txb_ft_r = _txb(slide, 10.0, footer_top + 0.03, 3.1, FOOTER_H - 0.04)
    _txt(txb_ft_r.text_frame,
         f"SLIDE {slide_idx}/{total_slides}  ·  DOKUMEN INTERNAL",
         7, bold=True, color=COLOR_RED_LIGHT, align=PP_ALIGN.RIGHT)



# ─────────────────────────────────────────────────────────────────────────────
# SLIDE TERAKHIR: Penutup
# ─────────────────────────────────────────────────────────────────────────────
def _slide_closing(prs, author: str, report_date: str, report_title: str,
                   total_slides: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background split: kiri merah, kanan abu muda
    lp = slide.shapes.add_shape(1, 0, 0, Inches(6.2), SLIDE_H)
    _fill(lp, COLOR_RED_DARK)

    rp = slide.shapes.add_shape(1, Inches(6.2), 0, Inches(7.13), SLIDE_H)
    _fill(rp, COLOR_GRAY_BG)

    # Garis vertikal emas
    sv = slide.shapes.add_shape(1, Inches(6.17), 0, Inches(0.07), SLIDE_H)
    _fill(sv, COLOR_GOLD)

    # Strip atas & bawah
    for y in (0, Inches(7.38)):
        s = slide.shapes.add_shape(1, 0, y, SLIDE_W, Inches(0.12))
        _fill(s, COLOR_RED_LIGHT)

    # Kiri: pesan penutup
    txb_logo = _txb(slide, 0.4, 0.35, 5.3, 0.75)
    _txt(txb_logo.text_frame, "TELKOMSEL", 30, bold=True,
         color=COLOR_WHITE, font="Poppins")

    txb_tg = _txb(slide, 0.4, 1.0, 5.5, 0.32)
    _txt(txb_tg.text_frame, "The Biggest & Most Innovative Digital Telco of Indonesia",
         8.5, italic=True, color=RGBColor(0xFF, 0xCC, 0xCC))

    gold_line = slide.shapes.add_shape(1, Inches(0.4), Inches(1.45), Inches(5.0), Inches(0.05))
    _fill(gold_line, COLOR_GOLD)

    txb_msg = _txb(slide, 0.4, 1.65, 5.4, 2.4)
    txb_msg.text_frame.word_wrap = True
    _txt(txb_msg.text_frame,
         "Terima kasih.\n\nLaporan ini disusun secara otomatis menggunakan "
         "Telkomsel Report Automator berdasarkan data tracking activity operasional.",
         15, color=COLOR_WHITE)

    txb_quote = _txb(slide, 0.4, 4.2, 5.4, 0.95)
    _txt(txb_quote.text_frame, '"Leading in Digital Connectivity"',
         14, italic=True, bold=True, color=COLOR_GOLD)

    txb_disc = _txb(slide, 0.4, 6.6, 5.4, 0.5)
    _txt(txb_disc.text_frame,
         "Dokumen ini bersifat RAHASIA. Hanya untuk keperluan internal PT. Telkomsel.",
         7.5, italic=True, color=RGBColor(0xFF, 0xCC, 0xCC))

    # Kanan: ringkasan dokumen
    txb_rt = _txb(slide, 6.5, 0.5, 6.5, 0.65)
    _txt(txb_rt.text_frame, "Ringkasan Dokumen", 18, bold=True,
         color=COLOR_RED_DARK)

    rl = slide.shapes.add_shape(1, Inches(6.5), Inches(1.1), Inches(6.4), Inches(0.05))
    _fill(rl, COLOR_RED)

    doc_info = [
        ("Judul Laporan",         report_title),
        ("Tanggal Dikeluarkan",   report_date),
        ("Disusun Oleh",          author),
        ("Diklasifikasikan",      "RAHASIA — INTERNAL"),
        ("Dibuat Dengan",         "Telkomsel Report Automator v3.0"),
        ("Total Slide",           f"{total_slides} Slide"),
    ]
    y_di = 1.3
    for key, val in doc_info:
        card = slide.shapes.add_shape(
            1, Inches(6.5), Inches(y_di), Inches(6.5), Inches(0.8)
        )
        _fill_border(card, COLOR_WHITE, RGBColor(0xE0, 0xE0, 0xE5), 0.5)

        txb_k = _txb(slide, 6.62, y_di + 0.07, 6.2, 0.24)
        _txt(txb_k.text_frame, key, 7.5, color=RGBColor(0x6B, 0x7B, 0x8D))

        txb_v = _txb(slide, 6.62, y_di + 0.33, 6.2, 0.35)
        _txt(txb_v.text_frame, val, 11.5, bold=True, color=COLOR_DARK_NAVY)

        y_di += 0.88


# ─────────────────────────────────────────────────────────────────────────────
# FUNGSI UTAMA: generate_pptx
# ─────────────────────────────────────────────────────────────────────────────
def generate_pptx(
    title: str,
    subtitle: str,
    author: str,
    report_date: str,
    df_raw: Optional[pd.DataFrame] = None,
    site_images: Optional[Dict[str, List[Dict]]] = None,
    # Parameter lama (backward-compat, tidak digunakan)
    bullets: Optional[List[str]] = None,
    chart_buf: Optional[bytes] = None,
    chart_title: str = "",
    analysis_text: str = "",
    x_col: str = "",
    y_col: str = "",
    user_images: Optional[List[dict]] = None,
) -> bytes:
    """
    Fungsi utama pembuatan PPTX Coverage Report Telkomsel.

    Parameter:
    ----------
    title       : Judul laporan (Slide Cover)
    subtitle    : Sub-judul / departemen (Slide Cover)
    author      : Nama pembuat
    report_date : Tanggal laporan (string)
    df_raw      : DataFrame dari Excel — setiap baris = 1 slide data
    site_images : Dict key=SITE_ID, value=List[Dict]
                  Setiap Dict: {"bytes": <bytes>, "label": <str>}
                  Urutan: [0]=Payload Surrounding, [1]=Maps/COVMO, [2..4]=Support images
    """
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    _patch_theme_font(prs, 'Poppins')   # ← ganti default theme font ke Poppins

    # Siapkan data baris
    rows_data: List[Dict] = []
    if df_raw is not None and not df_raw.empty:
        # Normalisasi nama kolom (strip + upper untuk matching fleksibel)
        df_norm = df_raw.copy()
        df_norm.columns = [str(c).strip() for c in df_norm.columns]
        rows_data = df_norm.to_dict(orient="records")

    total_data_slides = len(rows_data)
    total_slides = 1 + total_data_slides + 1  # cover + data + penutup

    # Slide 1: Cover
    _slide_1_cover(prs, title, subtitle, author, report_date)

    # Slide 2..N: Satu per baris data
    site_images = site_images or {}
    for i, row in enumerate(rows_data):
        slide_idx = i + 2  # slide ke-2 dst
        site_id   = str(row.get(COL_SITE_ID, f"SITE_{i+1}")).strip()
        images    = site_images.get(site_id, [])
        _slide_data_site(prs, row, images, slide_idx, total_slides)

    # Slide terakhir: Penutup
    _slide_closing(prs, author, report_date, title, total_slides)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    raw = buf.read()
    # Embed font Poppins ke dalam PPTX agar tampil benar di Canva / PowerPoint
    return embed_poppins_into_pptx(raw)
